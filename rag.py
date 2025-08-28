import os
import re
import time
import uuid
from typing import List, Dict, Any, Optional, Literal
import numpy as np
import torch
import google.generativeai as genai
from transformers import AutoTokenizer, AutoModel
import pinecone
import logging
from PyPDF2 import PdfReader
from pptx import Presentation
from dataclasses import dataclass
from datetime import datetime
import requests
import json
import psutil
from rouge_score import rouge_scorer
from bert_score import score as bert_score_calc
from nltk.translate.bleu_score import sentence_bleu
from nltk.translate.meteor_score import meteor_score as meteor_score_calc
from nltk import word_tokenize
import nltk

# Download NLTK data for tokenization
nltk.download('punkt')
nltk.download('wordnet')


# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Model provider types
ModelProvider = Literal["google", "ollama"]

@dataclass
class QueryMetrics:
    """Data class to store query-related metrics"""
    query: str
    words_in_answer: int
    query_processing_time: float
    total_latency: float
    timestamp: datetime
    top_k: int
    response_length: int
    llm_model: str
    cpu_percent: float
    ram_percent: float
    avg_match_score: float
    peak_memory_mb: float
    context_length: int
    embedding_time: float
    retrieval_time: float
    llm_time: float
    retrieval_quality: float
    # NLP evaluation metrics (requires ground truth)
    rouge_1: float = 0.0
    rouge_2: float = 0.0
    rouge_l: float = 0.0
    bert_score_precision: float = 0.0
    bert_score_recall: float = 0.0
    bert_score_f1: float = 0.0
    bleu_score: float = 0.0
    meteor_score: float = 0.0

@dataclass
class FileProcessingMetrics:
    """Data class to store file processing metrics"""
    filename: str
    file_type: str
    file_size_bytes: int
    processing_time: float
    text_length: int
    embedding_time: float
    storage_time: float
    timestamp: datetime
    embedding_model: str

class Metrics:
    """Class to handle metrics collection and reporting"""
    
    def __init__(self):
        self.query_metrics: List[QueryMetrics] = []
        self.file_metrics: List[FileProcessingMetrics] = []
    
    def add_query_metrics(self, metrics: QueryMetrics):
        """Add query metrics to the collection"""
        self.query_metrics.append(metrics)
        logger.info(f"Query metrics recorded: {metrics.words_in_answer} words, "
                   f"{metrics.query_processing_time:.2f}s processing, "
                   f"{metrics.total_latency:.2f}s total")
    
    def add_file_metrics(self, metrics: FileProcessingMetrics):
        """Add file processing metrics to the collection"""
        self.file_metrics.append(metrics)
        logger.info(f"File metrics recorded: {metrics.filename} processed in "
                   f"{metrics.processing_time:.2f}s")
    
    def get_query_stats(self) -> Dict[str, Any]:
        """Get statistics for all queries"""
        if not self.query_metrics:
            return {}
        
        # Existing metrics
        words_list = [m.words_in_answer for m in self.query_metrics]
        processing_times = [m.query_processing_time for m in self.query_metrics]
        total_latencies = [m.total_latency for m in self.query_metrics]
        cpu_percentages = [m.cpu_percent for m in self.query_metrics]
        ram_percentages = [m.ram_percent for m in self.query_metrics]
        match_scores = [m.avg_match_score for m in self.query_metrics]
        peak_memory_mbs = [m.peak_memory_mb for m in self.query_metrics]
        
        # New metrics
        context_lengths = [m.context_length for m in self.query_metrics]
        embedding_times = [m.embedding_time for m in self.query_metrics]
        retrieval_times = [m.retrieval_time for m in self.query_metrics]
        llm_times = [m.llm_time for m in self.query_metrics]
        retrieval_qualities = [m.retrieval_quality for m in self.query_metrics]
        
        # NLP metrics (assuming they are populated)
        rouge_1_scores = [m.rouge_1 for m in self.query_metrics]
        rouge_2_scores = [m.rouge_2 for m in self.query_metrics]
        rouge_l_scores = [m.rouge_l for m in self.query_metrics]
        bert_scores_p = [m.bert_score_precision for m in self.query_metrics]
        bert_scores_r = [m.bert_score_recall for m in self.query_metrics]
        bert_scores_f1 = [m.bert_score_f1 for m in self.query_metrics]
        bleu_scores = [m.bleu_score for m in self.query_metrics]
        meteor_scores = [m.meteor_score for m in self.query_metrics]
        
        # Calculate throughput (queries per second)
        total_time_all_queries = sum(total_latencies)
        throughput = len(self.query_metrics) / total_time_all_queries if total_time_all_queries > 0 else 0
        
        # Get unique LLM models used
        llm_models = list(set([m.llm_model for m in self.query_metrics]))
        
        return {
            "total_queries": len(self.query_metrics),
            "avg_words_in_answer": np.mean(words_list),
            "min_words_in_answer": min(words_list),
            "max_words_in_answer": max(words_list),
            "avg_query_processing_time": np.mean(processing_times),
            "avg_total_latency": np.mean(total_latencies),
            "fastest_query": min(processing_times),
            "slowest_query": max(processing_times),
            "llm_models_used": llm_models,
            "primary_llm_model": llm_models[0] if llm_models else "Unknown",
            "avg_cpu_percent": np.mean(cpu_percentages),
            "max_cpu_percent": max(cpu_percentages),
            "avg_ram_percent": np.mean(ram_percentages),
            "max_ram_percent": max(ram_percentages),
            "avg_match_score": np.mean(match_scores),
            "max_match_score": max(match_scores),
            "avg_peak_memory_mb": np.mean(peak_memory_mbs),
            "max_peak_memory_mb": max(peak_memory_mbs),
            
            # New stats
            "avg_context_length": np.mean(context_lengths),
            "avg_embedding_time": np.mean(embedding_times),
            "avg_retrieval_time": np.mean(retrieval_times),
            "avg_llm_time": np.mean(llm_times),
            "avg_retrieval_quality": np.mean(retrieval_qualities),
            "throughput": throughput,
            
            # NLP stats
            "avg_rouge_1": np.mean(rouge_1_scores),
            "avg_rouge_2": np.mean(rouge_2_scores),
            "avg_rouge_l": np.mean(rouge_l_scores),
            "avg_bert_score_precision": np.mean(bert_scores_p),
            "avg_bert_score_recall": np.mean(bert_scores_r),
            "avg_bert_score_f1": np.mean(bert_scores_f1),
            "avg_bleu_score": np.mean(bleu_scores),
            "avg_meteor_score": np.mean(meteor_scores),
        }
    
    def get_file_stats(self) -> Dict[str, Any]:
        """Get statistics for all file processing"""
        if not self.file_metrics:
            return {}
        
        processing_times = [m.processing_time for m in self.file_metrics]
        file_sizes = [m.file_size_bytes for m in self.file_metrics]
        
        # Get unique embedding models used
        embedding_models = list(set([m.embedding_model for m in self.file_metrics]))
        
        return {
            "total_files_processed": len(self.file_metrics),
            "avg_processing_time": sum(processing_times) / len(processing_times),
            "total_processing_time": sum(processing_times),
            "avg_file_size_mb": sum(file_sizes) / len(file_sizes) / (1024 * 1024),
            "fastest_file_processing": min(processing_times),
            "slowest_file_processing": max(processing_times),
            "embedding_models_used": embedding_models,
            "primary_embedding_model": embedding_models[0] if embedding_models else "Unknown"
        }
    
    def print_summary(self):
        """Print a summary of all metrics"""
        print("\n" + "="*50)
        print("METRICS SUMMARY")
        print("="*50)
        
        query_stats = self.get_query_stats()
        if query_stats:
            print("\nQUERY METRICS:")
            print(f"  Total queries: {query_stats['total_queries']}")
            print(f"  LLM Model: {query_stats['primary_llm_model']}")
            print(f"  Average words in answer: {query_stats['avg_words_in_answer']:.1f}")
            print(f"  Average query processing time: {query_stats['avg_query_processing_time']:.2f}s")
            print(f"  Average total latency: {query_stats['avg_total_latency']:.2f}s")
            print(f"  Throughput: {query_stats['throughput']:.2f} queries/sec")

            print("\nPERFORMANCE METRICS:")
            print(f"  Average embedding time: {query_stats['avg_embedding_time']:.2f}s")
            print(f"  Average retrieval time: {query_stats['avg_retrieval_time']:.2f}s")
            print(f"  Average LLM time: {query_stats['avg_llm_time']:.2f}s")
            
            print("\nRETRIEVAL METRICS:")
            print(f"  Average context length: {query_stats['avg_context_length']:.1f} chars")
            print(f"  Average retrieval quality (score): {query_stats['avg_retrieval_quality']:.3f}")

            print("\nNLP EVALUATION METRICS (requires ground truth):")
            print(f"  Average ROUGE-1: {query_stats['avg_rouge_1']:.3f}")
            print(f"  Average ROUGE-2: {query_stats['avg_rouge_2']:.3f}")
            print(f"  Average ROUGE-L: {query_stats['avg_rouge_l']:.3f}")
            print(f"  Average BERT Score (P): {query_stats['avg_bert_score_precision']:.3f}")
            print(f"  Average BERT Score (R): {query_stats['avg_bert_score_recall']:.3f}")
            print(f"  Average BERT Score (F1): {query_stats['avg_bert_score_f1']:.3f}")
            print(f"  Average BLEU Score: {query_stats['avg_bleu_score']:.3f}")
            print(f"  Average METEOR Score: {query_stats['avg_meteor_score']:.3f}")

            # System metrics
            print("\nSYSTEM METRICS:")
            print(f"  Average CPU usage: {query_stats['avg_cpu_percent']:.1f}%")
            print(f"  Peak CPU usage: {query_stats['max_cpu_percent']:.1f}%")
            print(f"  Average RAM usage: {query_stats['avg_ram_percent']:.1f}%")
            print(f"  Peak RAM usage: {query_stats['max_ram_percent']:.1f}%")
            print(f"  Average peak memory: {query_stats['avg_peak_memory_mb']:.1f} MB")
            print(f"  Maximum peak memory: {query_stats['max_peak_memory_mb']:.1f} MB")
            
            # Show all models if multiple are used
            if len(query_stats['llm_models_used']) > 1:
                print(f"  Models used: {', '.join(query_stats['llm_models_used'])}")
        
        file_stats = self.get_file_stats()
        if file_stats:
            print("\nFILE PROCESSING METRICS:")
            print(f"  Total files processed: {file_stats['total_files_processed']}")
            print(f"  Embedding Model: {file_stats['primary_embedding_model']}")
            print(f"  Average processing time: {file_stats['avg_processing_time']:.2f}s")
            print(f"  Total processing time: {file_stats['total_processing_time']:.2f}s")
            print(f"  Average file size: {file_stats['avg_file_size_mb']:.2f} MB")
            print(f"  Fastest file processing: {file_stats['fastest_file_processing']:.2f}s")
            print(f"  Slowest file processing: {file_stats['slowest_file_processing']:.2f}s")
            
            # Show all embedding models if multiple are used
            if len(file_stats['embedding_models_used']) > 1:
                print(f"  Embedding models used: {', '.join(file_stats['embedding_models_used'])}")
        
        print("="*50)

class OllamaClient:
    """Client for interacting with Ollama local models"""
    
    def __init__(self, base_url: str = "http://localhost:11434", model: str = "llama3.1:8b", timeout: int = 300, test_timeout: int = 10):
        self.base_url = base_url
        self.model = model
        self.api_url = f"{base_url}/api/generate"
        self.timeout = timeout
        self.test_timeout = test_timeout
    
    def generate_content(self, prompt: str):
        """Generate content using Ollama model"""
        try:
            payload = {
                "model": self.model,
                "prompt": prompt,
                "stream": False
            }
            
            response = requests.post(self.api_url, json=payload, timeout=self.timeout)
            response.raise_for_status()
            
            result = response.json()
            response_text = result.get("response", "")
            
            # Create a response object that matches Google's API
            class OllamaResponse:
                def __init__(self, text):
                    self.text = text
            
            return OllamaResponse(response_text)
            
        except requests.exceptions.RequestException as e:
            logger.error(f"Error calling Ollama API: {str(e)}")
            raise Exception(f"Failed to generate content with Ollama: {str(e)}")
        except Exception as e:
            logger.error(f"Error in Ollama content generation: {str(e)}")
            raise
    
    def test_connection(self) -> bool:
        """Test if Ollama is running and accessible"""
        try:
            response = requests.get(f"{self.base_url}/api/tags", timeout=self.test_timeout)
            return response.status_code == 200
        except:
            return False

class RAGSystem:
    def __init__(self, config: dict):
        self.pinecone_api_key = config["pinecone_api_key"]
        self.pinecone_environment = config.get("pinecone_environment", "us-east-1")
        self.pinecone_index_name = config.get("pinecone_index_name", "ibm-index")
        self.model_provider = config["model_provider"]
        self.llm_config = config["llm_config"]
        self.embedding_model_name = config["embedding_model_name"]
        self.index = None
        self.tokenizer = None
        self.model = None
        self.llm = None
        self.metrics = Metrics()
        self._init_pinecone()
        self._init_embedding_model()
        self._init_llm()
    
    def _init_pinecone(self):
        try:
            pc = pinecone.Pinecone(
                api_key=self.pinecone_api_key,
                environment=self.pinecone_environment
            )
            
            index_name = self.pinecone_index_name
            existing_indexes = pc.list_indexes().names()
            
            if index_name not in existing_indexes:
                logger.info(f"Creating new index: {index_name}")
                spec = pinecone.ServerlessSpec(
                    cloud='aws',
                    region=self.pinecone_environment
                )
                pc.create_index(
                    name=index_name,
                    dimension=768,
                    metric="cosine",
                    spec=spec
                )
                logger.info(f"Index {index_name} created successfully")
            
            self.index = pc.Index(index_name)
            
        except Exception as e:
            logger.error(f"Error initializing Pinecone: {str(e)}")
            raise
    
    def _init_embedding_model(self):
        try:
            logger.info("Initializing embedding model...")
            self.tokenizer = AutoTokenizer.from_pretrained(self.embedding_model_name, trust_remote_code=True)
            self.model = AutoModel.from_pretrained(self.embedding_model_name, trust_remote_code=True)
            self.model.eval()
        except Exception as e:
            logger.error(f"Error initializing embedding model: {str(e)}")
            raise
    
    def _init_llm(self):
        try:
            if self.model_provider == "google":
                api_key = self.llm_config["api_key"]
                model_name = self.llm_config["model"]
                if not api_key:
                    raise ValueError("Google API key is required when using Google models")
                genai.configure(api_key=api_key)
                self.llm = genai.GenerativeModel(model_name)
                self.llm_model_name = model_name
                logger.info(f"Initialized Google Gemini model: {model_name}")
                
            elif self.model_provider == "ollama":
                base_url = self.llm_config["base_url"]
                model_name = self.llm_config["model"]
                timeout = self.llm_config.get("timeout", 300)  # Default to 300s if not specified
                test_timeout = self.llm_config.get("test_timeout", 10)  # Default to 10s if not specified
                self.llm = OllamaClient(base_url, model_name, timeout, test_timeout)
                self.llm_model_name = model_name
                if not self.llm.test_connection():
                    raise Exception(f"Cannot connect to Ollama at {base_url}. Make sure Ollama is running.")
                logger.info(f"Initialized Ollama model: {model_name} (timeout: {timeout}s)")
                
            else:
                raise ValueError(f"Unsupported model provider: {self.model_provider}")
                
        except Exception as e:
            logger.error(f"Error initializing LLM: {str(e)}")
            raise
    
    def extract_text_from_file(self, file_path: str) -> str:
        start_time = time.time()
        try:
            if file_path.lower().endswith('.pdf'):
                text = self._extract_text_from_pdf(file_path)
            elif file_path.lower().endswith(('.pptx', '.ppt')):
                text = self._extract_text_from_ppt(file_path)
            else:
                raise ValueError(f"Unsupported file format: {file_path}")
            
            processing_time = time.time() - start_time
            
            # Record file processing metrics
            file_size = os.path.getsize(file_path)
            file_metrics = FileProcessingMetrics(
                filename=os.path.basename(file_path),
                file_type=os.path.splitext(file_path)[1].lower(),
                file_size_bytes=file_size,
                processing_time=processing_time,
                text_length=len(text),
                embedding_time=0.0,  # Will be updated in store_in_pinecone
                storage_time=0.0,    # Will be updated in store_in_pinecone
                timestamp=datetime.now(),
                embedding_model=self.embedding_model_name
            )
            self.metrics.add_file_metrics(file_metrics)
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    def _extract_text_from_pdf(self, file_path: str) -> str:
        text = []
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text.append(page.extract_text())
        return ' '.join(text)
    
    def _extract_text_from_ppt(self, file_path: str) -> str:
        text = []
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return ' '.join(text)
    
    def preprocess_text(self, text: str) -> str:
        try:
            # Convert to lowercase
            text = text.lower()
            # Remove special characters but keep periods and question marks
            text = re.sub(r'[^a-zA-Z\s\.\?]', '', text)
            # Remove extra whitespace
            text = ' '.join(text.split())
            return text
        except Exception as e:
            logger.error(f"Error preprocessing text: {str(e)}")
            raise
    
    def generate_embeddings(self, texts: List[str]) -> np.ndarray:
        start_time = time.time()
        try:
            embeddings = []
            for text in texts:
                inputs = self.tokenizer(text, padding=True, truncation=True, max_length=512, return_tensors="pt")
                with torch.no_grad():
                    outputs = self.model(**inputs)
                    embedding = outputs.last_hidden_state[:, 0, :].numpy()
                    embeddings.append(embedding[0])
            
            embedding_time = time.time() - start_time
            
            # Update the most recent file metrics with embedding time
            if self.metrics.file_metrics:
                self.metrics.file_metrics[-1].embedding_time = embedding_time
            
            return np.array(embeddings)
        except Exception as e:
            logger.error(f"Error generating embeddings: {str(e)}")
            raise
    
    def store_in_pinecone(self, texts: List[str], embeddings: np.ndarray, metadata: List[Dict] = None):
        start_time = time.time()
        try:
            batch_size = 100
            total_batches = (len(texts) + batch_size - 1) // batch_size
            
            for i in range(0, len(texts), batch_size):
                batch_texts = texts[i:i + batch_size]
                batch_embeddings = embeddings[i:i + batch_size]
                
                # Generate unique IDs using UUID and filenames if available
                ids = []
                for j in range(len(batch_texts)):
                    if metadata and metadata[j] and 'filename' in metadata[j]:
                        # Create ID based on filename and UUID to ensure uniqueness
                        file_id = f"{metadata[j]['filename'].replace('.', '_')}_{str(uuid.uuid4())[:8]}"
                        ids.append(file_id)
                    else:
                        # Fallback to UUID if no filename is available
                        ids.append(f"doc_{str(uuid.uuid4())}")
                
                vectors = list(zip(ids, 
                                 batch_embeddings.tolist(),
                                 [{"text": t, **(metadata[j] if metadata else {})} for j, t in enumerate(batch_texts)]))
                
                self.index.upsert(vectors=vectors)
                logger.info(f"Uploaded batch {(i // batch_size) + 1}/{total_batches} to Pinecone")
            
            storage_time = time.time() - start_time
            
            # Update the most recent file metrics with storage time
            if self.metrics.file_metrics:
                self.metrics.file_metrics[-1].storage_time = storage_time
                
        except Exception as e:
            logger.error(f"Error storing embeddings in Pinecone: {str(e)}")
            raise
    
    def query(self, query: str, top_k: int = 3, ground_truth_answer: str = None) -> str:
        total_start_time = time.time()
        
        try:
            # Generate query embedding
            embedding_start_time = time.time()
            query_embedding = self.generate_embeddings([query])[0]
            embedding_time = time.time() - embedding_start_time
            
            # Query Pinecone
            retrieval_start_time = time.time()
            results = self.index.query(
                vector=query_embedding.tolist(),
                top_k=top_k,
                include_metadata=True
            )
            retrieval_time = time.time() - retrieval_start_time
            
            # Extract relevant documents
            relevant_docs = [result['metadata']['text'] for result in results['matches']]
            context = " ".join(relevant_docs)
            context_length = len(context)
            
            # Generate response using LLM
            prompt = f"""
            Context: {context}
            
            Question: {query}
            
            Please provide a detailed and accurate answer based on the context provided above.
            It should be short and accurate.
            If the context doesn't contain enough information to answer the question,
            please indicate that.
            """
            
            llm_start_time = time.time()
            response = self.llm.generate_content(prompt)
            response_text = response.text
            llm_time = time.time() - llm_start_time
            
            # Calculate metrics
            total_latency = time.time() - total_start_time
            words_in_answer = len(response_text.split())
            avg_match_score = np.mean([m['score'] for m in results['matches'] if results['matches']]) if results['matches'] else 0.0

            # NLP evaluation metrics
            rouge_1, rouge_2, rouge_l = 0.0, 0.0, 0.0
            bert_p, bert_r, bert_f1 = 0.0, 0.0, 0.0
            bleu, meteor = 0.0, 0.0

            if ground_truth_answer:
                logger.info(f"Calculating NLP metrics with ground truth: {ground_truth_answer}")
                # logger.info(f"Generated response: {response_text}")

                # ROUGE
                scorer = rouge_scorer.RougeScorer(['rouge1', 'rouge2', 'rougeL'], use_stemmer=True)
                scores = scorer.score(ground_truth_answer, response_text)
                rouge_1 = scores['rouge1'].fmeasure
                rouge_2 = scores['rouge2'].fmeasure
                rouge_l = scores['rougeL'].fmeasure
                logger.info(f"ROUGE scores: 1={rouge_1:.3f}, 2={rouge_2:.3f}, L={rouge_l:.3f}")

                # BERTScore
                bert_p, bert_r, bert_f1 = bert_score_calc([response_text], [ground_truth_answer], lang="en", verbose=False)
                bert_p, bert_r, bert_f1 = bert_p.mean().item(), bert_r.mean().item(), bert_f1.mean().item()
                logger.info(f"BERT scores: P={bert_p:.3f}, R={bert_r:.3f}, F1={bert_f1:.3f}")

                # BLEU
                bleu = sentence_bleu([ground_truth_answer.split()], response_text.split())
                logger.info(f"BLEU score: {bleu:.3f}")

                # METEOR
                meteor = meteor_score_calc([word_tokenize(ground_truth_answer)], word_tokenize(response_text))
                logger.info(f"METEOR score: {meteor:.3f}")
            
            # Record query metrics
            query_metrics = QueryMetrics(
                query=query,
                words_in_answer=words_in_answer,
                query_processing_time=llm_time,  # Renamed for clarity
                total_latency=total_latency,
                timestamp=datetime.now(),
                top_k=top_k,
                response_length=len(response_text),
                llm_model=self.llm_model_name,
                cpu_percent=psutil.cpu_percent(),
                ram_percent=psutil.virtual_memory().percent,
                avg_match_score=avg_match_score,
                peak_memory_mb=psutil.virtual_memory().used / (1024 * 1024),
                context_length=context_length,
                embedding_time=embedding_time,
                retrieval_time=retrieval_time,
                llm_time=llm_time,
                retrieval_quality=avg_match_score,  # Using avg score as quality
                rouge_1=rouge_1,
                rouge_2=rouge_2,
                rouge_l=rouge_l,
                bert_score_precision=bert_p,
                bert_score_recall=bert_r,
                bert_score_f1=bert_f1,
                bleu_score=bleu,
                meteor_score=meteor
            )
            self.metrics.add_query_metrics(query_metrics)
            
            return response_text
            
        except Exception as e:
            logger.error(f"Error in query processing: {str(e)}")
            raise
    
    def get_metrics_summary(self):
        """Get a summary of all collected metrics"""
        return {
            "query_stats": self.metrics.get_query_stats(),
            "file_stats": self.metrics.get_file_stats()
        }
    
    def print_metrics_summary(self):
        """Print a formatted summary of all metrics"""
        self.metrics.print_summary() 